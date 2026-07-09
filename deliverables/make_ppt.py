"""Build the SalaryShield hackathon pitch deck (.pptx), one slide per scoring parameter."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG      = RGBColor(0x08, 0x0B, 0x11)
PANEL   = RGBColor(0x11, 0x16, 0x22)
CARD    = RGBColor(0x1A, 0x21, 0x33)
WHITE   = RGBColor(0xF9, 0xFA, 0xFB)
GREY    = RGBColor(0xB6, 0xBD, 0xCC)
MUTE    = RGBColor(0x8A, 0x92, 0xA3)
VIOLET  = RGBColor(0x8B, 0x5C, 0xF6)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
BLUE    = RGBColor(0x3B, 0x82, 0xF6)
DANGER  = RGBColor(0xEF, 0x44, 0x44)

FONT = "Segoe UI"
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=6):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (t, size, color, bold, italic) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
        r.font.italic = italic; r.font.name = FONT
    return tb


def bullets(slide, x, y, w, h, items, size=17, color=GREY, gap=10):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, (t, c) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0); p.alignment = PP_ALIGN.LEFT
        b = p.add_run(); b.text = "▸  "; b.font.size = Pt(size); b.font.color.rgb = VIOLET; b.font.bold = True; b.font.name = FONT
        r = p.add_run(); r.text = t; r.font.size = Pt(size); r.font.color.rgb = (c or color); r.font.name = FONT
    return tb


def header(slide, badge, title, accent=VIOLET):
    rect(slide, 0, 0, SW, Inches(1.55), PANEL)
    rect(slide, 0, 0, Inches(0.16), Inches(1.55), accent)
    # badge chip
    chip = rect(slide, Inches(0.55), Inches(0.28), Inches(3.4), Inches(0.42), CARD, line=accent, line_w=1.5,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    ctf = chip.text_frame; ctf.word_wrap = False
    cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text = badge; cr.font.size = Pt(12.5); cr.font.bold = True
    cr.font.color.rgb = accent; cr.font.name = FONT
    txt(slide, Inches(0.5), Inches(0.72), Inches(12.3), Inches(0.8),
        [(title, 30, WHITE, True, False)])


def content_slide(badge, title, items, accent=VIOLET, note=None):
    s = prs.slides.add_slide(BLANK); bg(s)
    header(s, badge, title, accent)
    bullets(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(4.9), items, size=18, gap=13)
    if note:
        nb = rect(s, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.7), CARD, line=accent, line_w=1.25,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        ntf = nb.text_frame; ntf.word_wrap = True; ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        np = ntf.paragraphs[0]; np.alignment = PP_ALIGN.LEFT
        r0 = np.add_run(); r0.text = "In this build:  "; r0.font.size = Pt(13); r0.font.bold = True; r0.font.color.rgb = accent; r0.font.name = FONT
        r1 = np.add_run(); r1.text = note; r1.font.size = Pt(13); r1.font.color.rgb = GREY; r1.font.name = FONT
    return s

# ---------------- Slide 1: Title ----------------
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, Inches(2.35), SW, Inches(0.06), VIOLET)
txt(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(1.2),
    [("SalaryShield", 60, WHITE, True, False)])
txt(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(1.0),
    [("AI-Driven, Inflation-Adaptive Compensation Intelligence", 26, VIOLET, True, False)])
txt(s, Inches(0.9), Inches(3.35), Inches(11.5), Inches(1.4),
    [("Quantify real-wage erosion. Justify fair corrections. Retain talent.", 18, GREY, False, True),
     ("Two-sided platform for employees and employers — powered by a City-Level Inflation Index (CLII) and a GenAI Copilot.", 15, MUTE, False, False)])
# tech chips
chips = ["React 19", "Vite 8", "Google Gemini AI", "CLII Engine", "Live Employer Console"]
cx = Inches(0.9)
for c in chips:
    w = Inches(1.9)
    ch = rect(s, cx, Inches(5.2), w, Inches(0.5), CARD, line=VIOLET, line_w=1.25, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = ch.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = c; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    cx += w + Inches(0.25)
txt(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.6),
    [("Hackathon Submission  ·  Scored across 10 parameters", 13, MUTE, False, False)])

# ---------------- Slide 2: Problem (1) ----------------
content_slide("SCORING PARAMETER 1  ·  PROBLEM DEFINITION",
    "Nominal salary is not real income",
    [("Inflation in India's tech hubs silently erodes purchasing power — Hyderabad 8.4%, Mumbai 9.1%, Bengaluru 7.9%.", None),
     ("National CPI (MOSPI) under-weights the urban tech-worker basket: IT-corridor rent, education, food delivery, commute.", None),
     ("Employees can't quantify their real-wage loss or defend a corrective hike with data.", None),
     ("Employers lose talent to inflation lag but lack city-level visibility — each exit costs ~₹37.5L in replacement.", DANGER),
     ("Net effect: mis-priced compensation, avoidable attrition, and unfair pay across geographies.", None)],
    accent=DANGER)

# ---------------- Slide 3: Solution overview ----------------
content_slide("SOLUTION OVERVIEW",
    "One platform, both sides of the pay table",
    [("Employee: personal inflation rate, real-wage score, 12-month savings-depletion forecast, break-even raise %, an investment allocator, and an AI negotiation coach.", None),
     ("Employer: live compensation budget simulator, attrition-risk ROI, city inflation heatmap, and a computed pay-fairness audit.", None),
     ("CLII Engine: a City-Level Inflation Index fusing rent, food, transport, utilities and education per city.", None),
     ("GenAI Copilot: Google Gemini, answering natural-language questions grounded in the user's own live numbers.", None),
     ("Every module reads one shared model — so the story stays consistent end to end.", EMERALD)],
    accent=VIOLET)

# ---------------- Slide 4: Architecture (2) ----------------
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "SCORING PARAMETER 2  ·  SOLUTION ARCHITECTURE", "Layered, model-driven architecture", EMERALD)
s.shapes.add_picture("C:/Users/ADMIN/kkl-workspace/salaryshield/deliverables/architecture.png",
                     Inches(0.55), Inches(1.7), width=Inches(12.23))

# ---------------- Slide 5: Engineering (3) ----------------
content_slide("SCORING PARAMETER 3  ·  ENGINEERING / BUILD QUALITY",
    "Clean, verifiable, production-shaped",
    [("React 19 + Vite 8 SPA; one component per module; lucide-react icons; a CSS-variable design system (glassmorphism).", None),
     ("Domain logic extracted into src/lib/inflation.js — a single source of truth reused by all 5 modules (DRY), incl. analyzeWorkforce() for the Employer Console.", EMERALD),
     ("Pure, deterministic, testable functions: analyzeBudget(), savingsSeries(), analyzeWorkforce(), break-even & personal-inflation engine.", None),
     ("Cross-module consistency: changing city/salary in the Budget Planner live-syncs the Employee Portal, Coach and Copilot within 1 second.", None),
     ("Production build passes — 1,780 modules, ~86 KB gzip. oxlint clean. localStorage persistence survives refresh.", None),
     ("Verified live in-browser during development (real DOM checks, zero console errors).", None)],
    accent=EMERALD,
    note="the whole app compiles and runs; the model layer is unit-test-ready pure functions.")

# ---------------- Slide 6: AI Appropriateness (4) ----------------
content_slide("SCORING PARAMETER 4  ·  AI APPROPRIATENESS",
    "AI only where it earns its place",
    [("The hard numbers — inflation, break-even, ROI — stay in deterministic code for accuracy and auditability.", None),
     ("The LLM does what it is uniquely good at: reasoning over the numbers and turning them into plain-language guidance.", None),
     ("Clear split — Engine computes, AI explains. No AI-for-AI's-sake; no hallucinated math.", VIOLET),
     ("Output is genuinely language-shaped: personalized advice and ready-to-send negotiation scripts.", None)],
    accent=VIOLET)

# ---------------- Slide 7: AI Implementation (5) ----------------
content_slide("SCORING PARAMETER 5  ·  AI IMPLEMENTATION",
    "A real, live LLM — grounded and multi-turn",
    [("GenAI Copilot calls Google Gemini 2.0 Flash directly — a real model call, not a scripted response.", EMERALD),
     ("RAG-style grounding: a system prompt injects the user's live budget, computed metrics, and the full CLII dataset for every city — answers cite the user's real ₹ figures, not model priors.", None),
     ("Multi-turn: full conversation history is passed on every call, so follow-up questions stay in context.", None),
     ("One data story across surfaces: Budget Planner → Negotiation Coach → Copilot all read src/lib/inflation.js.", EMERALD),
     ("Error-handled: a clear in-chat message (not a silent failure) if the API key or network is unavailable.", None)],
    accent=VIOLET,
    note="the Copilot is wired to a real Gemini API call (src/lib/llm.js) with the user's live personal-inflation, break-even and target-salary numbers as grounding context.")

# ---------------- Slide 8: AI Impact (6) ----------------
content_slide("SCORING PARAMETER 6  ·  AI IMPACT",
    "Invisible loss becomes a clear action",
    [("Employee: surfaces a hidden ~₹1.03L/yr loss and generates a copy-paste, data-backed pitch → higher raise success.", None),
     ("Employer: converts attrition risk into ROI — retain a key hire, avoid ~₹37.5L in replacement cost.", None),
     ("Personalized inflation (e.g. 7.7% vs a 8.4% city average) — an insight no generic salary calculator gives.", VIOLET),
     ("Data → decision in seconds, for both the individual and the organization.", None)],
    accent=AMBER)

# ---------------- Slide 9: Business Relevance (7) ----------------
content_slide("SCORING PARAMETER 7  ·  BUSINESS RELEVANCE",
    "A large, recurring, two-sided need",
    [("Addressable users: 5M+ Indian IT professionals and every enterprise running India delivery centers.", None),
     ("The pain is universal and recurring — it resurfaces at every appraisal cycle.", None),
     ("Maps directly to top HR priorities: retention, pay equity, and compensation-budget optimization.", None),
     ("Two-sided design means both the employee and the employer have a reason to adopt.", EMERALD)],
    accent=BLUE)

# ---------------- Slide 10: Benefits & Viability (8) ----------------
content_slide("SCORING PARAMETER 8  ·  BENEFITS & VIABILITY",
    "Clear model, low cost, real moat",
    [("Freemium employee tier drives acquisition; B2B SaaS (per-seat / per-audit) drives revenue from employers.", None),
     ("Data moat: the proprietary CLII index compounds in value as coverage and history grow.", VIOLET),
     ("Lean infrastructure — static SPA + serverless LLM calls — keeps unit economics healthy.", None),
     ("Viable path: MVP is built; remaining pieces (live feeds, HRMS integration, auth) are known and incremental.", EMERALD)],
    accent=EMERALD)

# ---------------- Slide 11: Working Demo (9) ----------------
content_slide("SCORING PARAMETER 9  ·  WORKING DEMO",
    "Live, interactive, end to end",
    [("A running MVP — not slideware. Four-minute flow:", WHITE),
     ("1  Budget Planner — enter salary, city & expenses → live personal inflation, savings-depletion chart, break-even raise; push rent up → deficit alert fires.", None),
     ("2  Negotiation Coach — the pitch auto-cites those exact numbers → one-click copy.", None),
     ("3  GenAI Copilot — ask 'Am I underpaid?' → a real Gemini call answers with the same live figures.", None),
     ("4  Employer Console — move the hike slider → attrition risk, ROI and fairness metrics recompute live.", None)],
    accent=AMBER)

# ---------------- Slide 12: Differentiation (10) ----------------
content_slide("SCORING PARAMETER 10  ·  DIFFERENTIATION",
    "Why this is not another salary calculator",
    [("Personal, city-level, expense-weighted inflation — not a one-size national average.", VIOLET),
     ("Two-sided: employee guidance AND employer workforce economics in one product.", None),
     ("AI output is an artifact you can act on — a negotiation script — not just a number.", None),
     ("A single, consistent data story flows across every module.", EMERALD),
     ("The CLII index is defensible IP that strengthens over time.", None)],
    accent=VIOLET)

# ---------------- Slide 13: Roadmap / Close ----------------
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "ROADMAP & ASK", "From MVP to production", VIOLET)
bullets(s, Inches(0.7), Inches(1.95), Inches(6.0), Inches(4.6),
    [("Now: MVP live — CLII model, 5 modules, live Employer Console, budget engine, Gemini-powered Copilot.", EMERALD),
     ("Next: live data ingestion (MOSPI/RBI/rent/basket) feeding the CLII pipeline.", None),
     ("Next: server-side LLM proxy (key security, rate limiting) + response streaming.", None),
     ("Then: HRMS integration, auth, multi-user, historical CLII moat.", None)], size=17, gap=14)
card_box = rect(s, Inches(7.0), Inches(1.95), Inches(5.6), Inches(4.3), PANEL, line=VIOLET, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(7.3), Inches(2.2), Inches(5.0), Inches(3.8),
    [("The one-line pitch", 18, VIOLET, True, False),
     ("SalaryShield turns invisible inflation into a fair, data-backed pay decision — for the employee and the employer.", 17, WHITE, False, False),
     ("", 8, WHITE, False, False),
     ("Built with React 19 + Vite 8, a deterministic CLII engine, and a Gemini-powered GenAI Copilot grounded in the user's real numbers.", 14, GREY, False, False)], sp_after=10)

prs.save("C:/Users/ADMIN/kkl-workspace/salaryshield/deliverables/SalaryShield_Pitch_Deck.pptx")
print("saved SalaryShield_Pitch_Deck.pptx  ·  slides:", len(prs.slides._sldIdLst))
